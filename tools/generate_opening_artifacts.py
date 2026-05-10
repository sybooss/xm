from __future__ import annotations

from copy import deepcopy
import re
import shutil
import subprocess
import textwrap
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK, WD_TAB_ALIGNMENT, WD_TAB_LEADER
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Mm, Pt, RGBColor
import fitz
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_DOC = ROOT / "人工智能学院复杂软件系统实践课开题报告格式.doc"
TASKBOOK_PPTX = ROOT / "《复杂软件系统实践》任务书.pptx"

DOC_OUT = ROOT / "output" / "doc"
TMP_DIR = ROOT / "tmp" / "docs"
FIG_DIR = TMP_DIR / "figures"

FINAL_DOCX = DOC_OUT / "开题报告.docx"
FINAL_PDF = DOC_OUT / "开题报告.pdf"

TMP_TEMPLATE_DOC = TMP_DIR / "template_source.doc"
TMP_TEMPLATE_DOCX = TMP_DIR / "template_source.docx"
TMP_REPORT_DOCX = TMP_DIR / "opening_report.docx"
TMP_PASS1_DOCX = TMP_DIR / "opening_report_pass1.docx"

BODY_PAGE_START = 1
FIGURE_BG_THRESHOLD = 245
FIGURE_CROP_PADDING = 24

ARCH_FIG = FIG_DIR / "architecture.png"
FLOW_FIG = FIG_DIR / "tech_flow.png"
GPT_IMG_DIR = Path(r"C:\Users\86188\.codex\generated_images\019d96f0-76b3-74b1-9901-eb1a34207209")
GPT_ARCH_FIG = GPT_IMG_DIR / "ig_05fe9e3054033b9a0169e8a6e576ec81919f2b6ca23f44e95e.png"
GPT_FLOW_FIG = GPT_IMG_DIR / "ig_05fe9e3054033b9a0169e8be568c348191bf06ee85c2011c37.png"

TITLE = "基于检索增强的电商退换货智能客服系统设计与实现"
ACADEMY = "人工智能学院"
MAJOR = "软件工程"
CLASS_NAME = "软件231"
STUDENT_ID = "2025467065902"
STUDENT_NAME = "刘剑宇"
TEACHER = "王向东"
OPENING_DATE = "2026年4月27日"

BLACK = RGBColor(0, 0, 0)
BLUE = RGBColor(52, 87, 140)
LIGHT_BLUE = RGBColor(223, 233, 247)
LIGHT_GRAY = RGBColor(242, 242, 242)

FONT_CANDIDATES = [
    Path(r"C:\Windows\Fonts\msyh.ttc"),
    Path(r"C:\Windows\Fonts\simsun.ttc"),
    Path(r"C:\Windows\Fonts\simhei.ttf"),
]

SECTION_TITLES = {
    "background": "一、选题的背景、意义",
    "content": "二、研究的基本内容与拟解决的主要问题",
    "method": "三、研究的方法与技术路线、研究难点，预期达到的目标",
    "schedule": "四、论文详细工作进度和安排",
    "refs": "五、主要参考文献",
}

TOC_ENTRY_SPECS = [
    (1, SECTION_TITLES["background"], SECTION_TITLES["background"]),
    (2, "1.1 选题背景", "1.1 选题背景"),
    (2, "1.2 研究意义", "1.2 研究意义"),
    (2, "1.3 国内外研究现状", "1.3 国内外研究现状"),
    (2, "1.4 发展趋势", "1.4 发展趋势"),
    (1, SECTION_TITLES["content"], SECTION_TITLES["content"]),
    (2, "2.1 研究的基本内容", "2.1 研究的基本内容"),
    (2, "2.2 拟解决的主要问题", "2.2 拟解决的主要问题"),
    (1, SECTION_TITLES["method"], SECTION_TITLES["method"]),
    (2, "3.1 研究方法", "3.1 研究方法"),
    (2, "3.2 技术路线", "3.2 技术路线"),
    (2, "3.3 研究难点", "3.3 研究难点"),
    (2, "3.4 预期达到的目标", "3.4 预期达到的目标"),
    (1, SECTION_TITLES["schedule"], SECTION_TITLES["schedule"]),
    (1, SECTION_TITLES["refs"], SECTION_TITLES["refs"]),
]

BACKGROUND_PARAGRAPHS = [
    "随着电子商务平台和在线零售模式的快速普及，商品交易规模、订单处理规模和售后服务场景持续扩大。与传统线下售后相比，线上退货、换货、退款到账、物流异常、人工投诉等问题具有咨询量大、规则复杂、处理环节多、时效要求高等特点。用户在提出售后问题时，通常不仅希望得到“能不能办理”的回答，还希望同时知道规则依据、所需材料、办理时限以及下一步操作路径。因此，面向电商退换货场景构建智能客服系统，既具有明确的现实需求，也具有较强的工程落地价值。",
    "从平台经营机制看，退换货服务已不再只是订单完成后的补充环节，而是直接影响用户信任、复购意愿和平台品牌口碑的关键服务节点。特别是在直播电商、即时零售、跨境电商和多商家平台环境下，商品来源、物流责任、平台政策、商家规则与消费者权益之间形成了更加复杂的约束关系。人工客服虽然能够在一定程度上处理例外情况，但面对高并发咨询、规则频繁更新和标准口径要求时，往往存在响应不一致、培训成本高和服务效率波动明显等问题。",
    "近年来，大语言模型、检索增强生成、任务型对话系统等技术快速发展，推动了智能客服从关键词匹配式问答向多轮对话、知识增强问答和场景化交互演进。相比早期静态 FAQ 机器人，融合知识库检索和上下文理解的系统能够更好地处理复杂售后问题。然而，电商退换货场景带有明显的业务规则约束，不同商品类别、下单时间、订单状态、签收情况、赠品处理规则和平台政策都会影响最终答复。若完全依赖大模型自由生成，容易出现规则依据不足、答案不稳定甚至与平台规则不一致的问题，这与课程项目对“可验证、可解释、可展示”的要求并不完全吻合。",
    "同时，电商售后知识具有强时效、强条件组合和强文本解释性的特点。许多规则并不是单一条款即可直接给出结论，而是需要结合“是否签收”“是否超时”“是否影响二次销售”“是否属于质量问题”“是否已发起工单”等上下文条件综合判断。由此可见，真正有使用价值的智能客服系统，不应只是生成自然语言答案，更应具备规则约束、依据返回、过程可追溯和异常场景回退等能力。",
    "从用户体验角度出发，售后咨询并不是孤立的一问一答过程，而往往伴随着重复追问、确认条件、补充订单号、询问时限和对处理结果的不确定性。例如，用户可能先询问“这个商品能不能退”，随后继续追问“那运费谁承担”“多久到账”“如果商家拒绝怎么办”。如果系统缺乏上下文继承能力，就很难在后续轮次中保持一致口径，也无法体现出真正的服务连续性。因此，构建具备多轮记忆与业务状态跟踪能力的客服系统，已经成为当前场景化智能服务的重要要求。",
    "此外，课程项目与企业真实应用虽不完全等同，但两者在系统完整性要求上具有相通之处。一个能够支撑开题答辩、阶段验收和结项展示的项目，应当不仅能说清楚算法思路，更要能展示清晰的模块划分、稳定的接口调用、规范的数据组织和可验证的运行结果。因此，以电商退换货智能客服为课题对象，既有利于结合现实业务问题展开设计，也适合在课程要求下完成从需求分析到系统实现的全流程训练。",
    "从复杂软件系统实践课程的角度看，电商退换货智能客服系统并不是一个单点算法实验，而是一个包含前端交互、后端业务服务、数据存储、知识库管理、检索增强、对话状态维护和测试验证的综合软件系统。它既可以体现完整的软件工程链路，也便于在开题、开发、验收和结项阶段形成连续成果。围绕这一目标，本课题采用前后端分离和 AI 增强协同的思路，将 Vue 3 作为前端框架、Spring Boot 作为后端业务核心、MySQL 作为主要数据存储，并以 LangChain4j 作为 AI 增强层接入大模型与检索链路，从而在保证业务控制力的同时提升系统智能化程度。",
    "基于以上背景，选择“基于检索增强的电商退换货智能客服系统设计与实现”作为课题，一方面能够围绕真实业务问题展开系统设计，另一方面也能将需求分析、接口设计、数据库建模、知识组织、智能增强、系统测试和可视化展示等内容有机结合，使课题既具有较明确的应用场景，又具有适合课程考核的复杂软件系统特征。",
]

SIGNIFICANCE_PARAGRAPHS = [
    "本课题的研究意义主要可以从实际应用、工程训练、知识管理和后续扩展几个方面来理解。",
    "先从应用场景来看，电商平台里很多售后咨询都集中在退货条件、退款时间、物流异常、商家超时未处理这些问题上。这些问题看起来都不算特别复杂，但真正回答时往往要结合订单状态、平台规则和办理流程来判断。如果能把这部分高频问题先标准化，再用系统自动完成一部分说明和引导，就能减轻人工客服压力，也能让用户更快拿到清楚的处理建议。",
    "再从课程实践来看，这个课题比较适合做成一个完整的软件系统，而不是只做一个算法小实验。前端需要负责咨询页面和管理页面，后端要处理接口、规则判断和会话管理，数据库要保存订单、知识文档和日志，LangChain4j 负责把检索增强和大模型能力接进来。整个过程会涉及需求分析、数据库设计、接口联调、前后端配合和测试验收，和课程想训练的能力比较贴合。",
    "从知识管理角度看，很多客服规则平时都散落在帮助文档、培训话术或者人工经验里，时间一长就容易出现口径不统一、更新不同步的问题。本课题把 FAQ、平台规则和售后说明整理成可维护的知识库，再把检索结果和回答依据一起返回，这样后面不管是补充规则、修改政策，还是排查回答问题，都会更方便一些。",
    "从技术方案选择上看，这个课题也有一定的探索意义。现在很多系统喜欢把大模型直接放到最前面，但对于退换货这类规则比较明确的场景来说，完全交给模型自由生成其实风险不小。比较稳妥的做法还是让 Spring Boot 负责业务主干和规则控制，再把 LangChain4j 放在增强层，去做语义理解、提示词编排和检索增强。这样既能用上大模型的表达能力，也能尽量保证结果可控。",
    "另外，这个课题还有比较好的延续性。开题阶段主要是把选题、系统结构和实现路线说清楚，后面还能继续往知识文档扩充、提示词优化、日志分析、上下文记忆和系统测试这些方向深入。这样做出来的成果不只是开题能用，结项阶段也能继续沿着同一条路线往下完善。",
    "所以我觉得，这个课题的价值不在于做一个会聊天的页面，而在于把业务规则、知识库、多轮问答和后台管理真正整合成一个能运行、能展示、也方便后续完善的系统。这一点和复杂软件系统实践课程的要求是比较一致的。",
]

STATE_OF_ART_PARAGRAPHS = [
    "从已有研究来看，国外在任务型对话、知识增强问答和检索增强生成方面起步更早。Gao 等对会话式人工智能中的神经方法做了较系统的梳理，认为多轮对话、知识接入和生成控制是对话系统持续发展的几个关键方向[7]。Louvan 和 Magnini 也总结了任务型对话中的意图识别和槽位填充方法，说明对用户问题的理解能力一直都是系统效果的核心基础[10]。",
    "在检索增强方向上，Lewis 等提出的 RAG 框架把参数知识和外部知识检索结合起来，使回答结果在事实性和可追溯性方面得到提升[8]。Karpukhin 等提出的 DPR、Guu 等提出的 REALM，则进一步说明检索机制在复杂问答任务中确实能够带来明显帮助[9][11]。这些研究对本课题有直接参考意义，因为本课题同样需要解决“先找到合适的规则和知识，再组织成回答”这个问题。",
    "和开放领域聊天不同，电商客服更强调业务约束。JDDC 和 CrossWOZ 等数据集为中文多轮对话研究提供了数据基础[12][14]，Deng 等对电商商品问答的综述也指出，电商问答往往具有规则来源多、用户表达口语化、答案要求可解释等特点[13]。这说明电商售后问答不能只追求语言自然，还要把规则依据和处理流程放进去。",
    "国内在任务型对话和客服理解方面也积累了不少成果。曹亚如等对多轮任务型对话系统进行了综述，从系统结构、评估方法和应用场景等角度总结了模块式方法与端到端方法的差别[1]。黄健、李锋提出的 SPD-BERT 面向口语对话理解，引入角色信息和轮次结构后，对智能客服场景中的意图识别和实体识别有明显帮助[2]。",
    "在工程化研究上，周超等提出了面向工业运维对话的意图和语义槽联合识别方法，通过联合建模提升复杂场景下的理解效果[3]。任元凯、谢振平则针对大语言模型在专业场景中意图不够稳定的问题，提出了结合领域知识进行增强的思路[4]。这些工作说明，在专业场景里，单靠通用模型往往还不够，还是需要和领域规则、业务知识结合起来用。",
    "近两年，国内关于检索增强生成和知识增强大模型的综述也越来越多。李子骏等对 RAG 技术的核心流程进行了梳理[5]，曹荣荣等从知识图谱融合的角度分析了缓解幻觉和提升可解释性的做法[6]。从这些研究可以看出，当前的一个共同认识是：如果想把智能客服真正落到具体行业场景里，就不能只看生成能力，还要同时考虑检索、规则、上下文和系统控制。",
    "不过，现有研究里很多工作还是更偏向单点突破，比如专门研究意图识别、专门优化检索，或者只关注回答质量。对于课程项目来说，除了这些技术点，还必须考虑页面展示、接口联调、日志记录和现场演示稳定性。因此，本课题在借鉴已有研究的同时，也要结合实际开发条件做工程化整合，这样才更符合课程项目的要求。",
]

TREND_PARAGRAPHS = [
    "结合前面的研究现状，可以看出领域智能客服的发展大致有几个比较明确的方向。第一是从单轮问答走向多轮对话，系统不再只回答一句话，而是要能记住上一轮说了什么。第二是从单纯生成走向“规则 + 检索 + 生成”结合的模式，回答不仅要自然，还要有依据。第三是从通用聊天走向垂直行业场景，越是规则明确、流程固定的业务，越需要做领域化设计。",
    "从技术细节上看，检索增强系统也不再停留在简单的文档召回和拼接阶段，而是开始更重视查询改写、上下文记忆、结构化知识融合和结果校验。也就是说，系统不仅要把文档找出来，还要知道哪些内容该优先使用，哪些规则必须先判断，最后怎样把结果组织成用户容易理解的话。",
    "从工程建设角度看，智能客服正在从一个回答模块逐步变成一个完整的平台能力。除了问答本身，知识录入、版本管理、日志追踪、测试回放、权限控制和后台维护都越来越重要。对课程项目来说，这一点也很关键，因为答辩时不能只展示模型回答效果，还要把系统整体结构和维护能力讲清楚。",
    "基于这些趋势，本课题最后选择“前后端分离的业务系统 + LangChain4j AI 增强层”这条路线。这样做一方面比较符合当前智能客服的工程化方向，另一方面也更适合课程项目对完整性、稳定性和可展示性的要求。",
]

CONTENT_INTRO_PARAGRAPHS = [
    "围绕“基于检索增强的电商退换货智能客服系统设计与实现”这一课题，本文准备从业务需求、系统架构、知识管理、智能增强和系统展示五个方面展开研究。和把大模型直接当成项目核心的做法不同，本课题会把业务规则和系统控制权放在 Spring Boot 服务层，再通过 LangChain4j 提供语义理解、提示词编排和检索增强能力。",
    "具体来说，首先要围绕退货申请、换货申请、退款到账、物流异常、规则咨询和投诉转接这些典型场景，梳理用户常见提问和业务流程之间的对应关系；其次要整理平台规则、FAQ 文档和售后说明，构建可维护的知识库；最后再把多轮对话管理、订单上下文读取和回答生成链路接起来，让系统能够结合历史提问和订单信息给出更贴近真实业务的答复。",
    "除了核心问答功能，本课题还会把知识库维护、接口测试、日志记录和界面展示一起考虑进去。这样从开题开始，整个项目就不是零散的几个功能点，而是一个比较完整、后面也方便继续完善的系统原型。",
]

PROBLEM_PARAGRAPHS = [
    "本课题拟重点解决以下四个主要问题。",
    "第一，用户提问方式很灵活，怎么把售后意图分准。电商用户经常会说“这单能不能退”“寄回去多久能到账”“物流卡住了怎么办”这类比较口语化的话，字面上不长，但背后对应的是明确的业务流程。如果意图分错，后面的规则检索和回答生成基本都会跟着出问题，所以这一块是整个链路的起点。",
    "第二，平台规则比较复杂，怎么让回答既自然又有依据。退换货问题往往要同时看平台政策、商品状态、订单节点和是否超时等条件，用户关心的不只是结论，还会追问为什么、依据是什么。因此系统不能只给一个简短回答，还要把 FAQ、规则文档和订单信息结合起来。",
    "第三，连续追问很多，怎么保证上下文不乱。实际对话里，用户第二轮、第三轮通常不会把背景重新说一遍，而是直接问“那多久到账”“那商家不处理怎么办”。如果系统不能继承前文里的订单状态和场景信息，回答就很容易前后不一致，所以多轮上下文管理也是必须解决的问题。",
    "第四，AI 增强能力怎么和传统业务系统配合，而不是把原来的业务逻辑冲掉。如果完全交给 AI 去主导业务流程，规则容易跑偏；但如果一点 AI 能力都不用，系统在语义理解和表达上又会比较生硬。所以本课题需要处理好的一个核心问题，就是让 LangChain4j 做增强层，而让 Spring Boot 继续承担业务主干。",
    "这四个问题其实是连在一起的。意图识别影响后面的检索，检索结果又会影响回答生成，而上下文管理和异常回退决定了系统在真实使用和答辩展示时稳不稳。因此本课题不会把这些问题拆开单独处理，而是放在一条完整链路里一起设计。",
]

METHOD_PARAGRAPHS = [
    "本课题准备采用“文献调研 + 需求分析 + 原型设计 + 分层实现 + 测试验证”的方式推进。",
    "第一步是先把相关资料看清楚。主要会阅读任务型对话、检索增强生成、意图识别、电商问答和软件架构设计方面的文献，弄清楚这个方向现在常见的做法是什么、容易出问题的地方在哪里、哪些思路适合放到课程项目里使用。",
    "第二步是在需求分析的基础上做系统原型。这里不会一开始就去追求复杂算法，而是先把页面、接口、知识文档、订单数据、会话状态和 AI 调用边界这些内容分清楚，明确每个模块负责什么，先把整体框架搭起来。",
    "第三步是按前后端分离的方式逐步实现系统。前端用 Vue 3 搭建咨询工作台和知识库管理界面，后端用 Spring Boot 负责接口服务、业务规则判断和会话管理，MySQL 存储订单、知识文档、会话记录和日志信息，LangChain4j 作为增强层接入大模型和检索链路，负责提示词编排、上下文记忆和回答组织。",
    "第四步是围绕典型售后场景做测试和修正。比如退货申请、换货申请、退款进度、物流异常、规则咨询和投诉转接这些场景，都要准备对应的测试样例，看系统能不能稳定识别意图、找到合适规则，并给出比较完整的答复。",
    "在验证方式上，本文准备把案例测试和指标观察结合起来使用。一方面，通过一组固定案例去看系统在不同订单状态、不同规则条件下的回答表现；另一方面，再从意图识别是否准确、规则命中是否合理、回答是否完整、依据是否能追溯以及整体运行是否稳定这些角度做综合判断。",
    "在具体实施顺序上，本课题还是按“先规则、后增强；先单轮、后多轮；先原型、后优化”的思路推进。先把基础接口、规则问答和知识维护做稳定，再接入 LangChain4j 的检索增强和上下文记忆，最后根据测试结果继续调整提示模板、知识组织方式和异常处理逻辑。这样做比较稳，也更方便在答辩时说明系统是怎么一步一步完善起来的。",
]

TECH_ROUTE_PARAGRAPHS = [
    "在技术实现上，本课题采用前后端分离的整体架构，同时把 LangChain4j 明确放在 AI 增强层，而不是把它当成整个项目的核心。具体来说，Vue 3 负责咨询界面和知识库管理界面，Spring Boot 负责接口调度、业务规则处理、会话管理和订单信息读取，MySQL 用来保存订单数据、FAQ 文档、会话记录和系统日志，LangChain4j 则嵌入后端服务中，负责组织提示词、维护上下文和调用检索增强流程。",
    "系统运行时，用户先在前端输入售后问题，请求提交到 Spring Boot 接口层。后端先做参数校验和业务路由，再结合订单状态、知识文档和历史会话判断当前问题属于哪一类场景。在这个基础上，LangChain4j 再把检索到的规则、订单上下文和提示模板组织起来，调用大模型生成更自然的答复。",
    "考虑到课程答辩对稳定性的要求，系统不能完全依赖 AI 链路。如果大模型调用异常，或者增强链路暂时不可用，后端仍然要能根据已有规则给出基础答复，至少把能不能退、怎么处理、需要补充什么信息这些关键内容说明白。这样系统在演示时会更稳，回答也更可控。",
    "除了问答链路，本课题还会同时建设知识库管理和日志记录链路。也就是说，系统不仅要能回答问题，还要支持规则文档的录入、编辑、删除、分类查询和内容维护，并把命中规则、检索结果和回答输出记录下来，方便后面做调试、回放和优化。这部分内容也是复杂软件系统里不能缺的一块。",
    "整体来看，这条技术路线始终坚持“业务系统为主，AI 能力为辅”的思路。这样既能突出 Spring Boot 在业务控制上的核心作用，也能把 LangChain4j 的增强能力合理接进去，更适合本课题这种规则明确、需要稳定展示的课程项目。",
]

DIFFICULTY_PARAGRAPHS = [
    "本课题的研究难点主要体现在以下三个方面。",
    "第一，相似售后问题的意图区分比较难。现实里很多问题都不是特别标准，例如“退货多久到账”既可能是在问退货申请，也可能是在问退款进度；“物流不动了能不能投诉”又同时带有物流异常和人工转接的意思。怎么在这种表达交叉的情况下找准主意图，是系统设计里最容易出偏差的地方。",
    "第二，知识检索、业务规则和自然语言生成之间要做好平衡。系统既不能只会机械地贴规则，也不能为了回答自然就把规则说偏。要想把这两头都兼顾起来，就需要 Spring Boot 的业务判断、MySQL 里的知识文档以及 LangChain4j 的提示编排和生成过程配合稳定。",
    "第三，多轮上下文承接和系统展示稳定性需要同时考虑。用户连续追问时经常会省略背景，如果上下文管理不好，错误就会在后面几轮里不断放大。另外课程答辩对系统稳定性要求比较高，所以还要提前把 AI 调用异常、规则回退和兜底回答这些情况考虑进去。",
    "归纳来看，这几个难点最后都落在系统协同上，也就是意图识别、知识检索、订单上下文、回答生成和规则回退能不能在同一条链路里稳定配合。和单独做一个算法模块相比，这种问题更考验模块边界设计和工程实现质量。",
]

GOAL_PARAGRAPHS = [
    "本课题的预期目标如下。",
    "第一，完成一个面向电商退换货场景的智能客服系统原型，做到前后端能够联动，核心数据能够闭环。",
    "第二，前端基于 Vue 3 实现咨询工作台和知识库管理页面，后端基于 Spring Boot 提供统一接口、业务规则处理和会话管理服务。",
    "第三，系统基于 MySQL 完成订单数据、知识文档、会话记录和系统日志的结构化存储，支持 FAQ 和平台规则的维护、查询与检索。",
    "第四，将 LangChain4j 作为 AI 增强层接入系统，实现提示词编排、上下文记忆、知识增强问答和自然语言答复组织，但不改变 Spring Boot 作为业务核心的定位。",
    "第五，系统能够覆盖退货申请、换货申请、退款进度、物流异常、规则说明和投诉转接等典型售后场景，并在回答中体现依据和处理建议。",
    "第六，形成测试用例、展示脚本和答辩材料，使项目能够支撑开题、开发、中期检查和结项展示几个阶段的使用。",
    "在可验收指标上，本文计划让系统内置 30 条以上规则或 FAQ 文档，支持 7 类典型客服意图，能够完成多轮上下文问答，并实现知识文档的增删改查和后台展示。把目标写得更具体，也方便后面按这个标准去开发和验收。",
]

SCHEDULE_TAIL_PARAGRAPHS = [
    "开题答辩通过后，后续工作将按第10周至第12周的课程节奏推进系统实现、功能联调、测试验证与结项准备，确保在结项答辩阶段能够展示较完整的系统原型和支撑材料。",
]

TEAM_WORK_INTRO_PARAGRAPHS = [
    "本人刘剑宇担任本组组长，主要负责项目选题统筹、任务协调、进度推进和开题汇报组织，同时承担意图识别与多轮对话管理模块的设计与说明工作。该部分是系统从用户输入进入业务处理链路的起点，重点解决用户需求理解、关键信息提取和上下文连续承接问题。",
]

TEAM_WORK_TABLE = [
    ["身份", "负责模块", "具体完成工作", "材料撰写与答辩内容"],
    [
        "组长：刘剑宇",
        "意图识别与多轮对话管理模块",
        "负责梳理退货申请、换货申请、退款进度、物流异常、规则咨询和投诉转接等典型问题的意图类别；设计用户问题关键信息提取方式；规划会话状态保存、多轮追问和上下文承接逻辑；协调知识库检索、回答生成和前端展示之间的接口衔接，保证该模块能够支撑后续系统联调与演示。",
        "开题报告中主要撰写研究内容、拟解决的主要问题、本人负责模块说明和系统整体推进安排；PPT 中重点讲解意图识别流程、多轮对话示例、上下文管理思路以及该模块和知识库检索、回答生成之间的协作关系；开题汇报控制在 8 分钟内，突出本人作为组长负责的统筹协调和核心模块设计工作。",
    ],
]

RESEARCH_TABLE = [
    ["研究内容", "对应系统模块", "主要实现要点", "预期产出"],
    ["售后意图识别", "Spring Boot 业务服务层 + LangChain4j 意图分析组件", "对用户问题进行七类售后意图判断，并结合上下文提高识别准确度", "形成可用于后续检索与回复生成的意图标签和处理入口"],
    ["退换货知识库构建", "MySQL 知识文档表 + 知识管理后台", "整理 FAQ、平台规则、退款说明和异常处理方案，并支持维护更新", "形成结构化、可查询、可扩展的规则知识库"],
    ["多轮对话管理", "Vue 3 咨询工作台 + Spring Boot 会话管理服务", "维护会话历史、追问上下文和用户当前咨询状态", "实现连续问答和上下文承接能力"],
    ["订单上下文融合", "Spring Boot 订单服务 + 提示词装配模块", "读取订单号、商品信息、订单状态和时间节点，并参与答复生成", "输出更贴近真实业务流程的个性化答复"],
    ["AI 增强回复与系统展示", "LangChain4j 编排层 + Vue 3 管理后台", "组织检索增强、大模型调用、依据展示和文档增删改查", "形成可展示、可维护、可验收的完整系统原型"],
]

REFERENCES = [
    "[1] 曹亚如, 张丽萍, 赵乐乐. 多轮任务型对话系统研究进展[J]. 计算机应用研究, 2022, 39(2): 331-341.",
    "[2] 黄健, 李锋. 融合角色、结构和语义的口语对话预训练语言模型[J]. 计算机应用研究, 2022, 39(8): 2397-2402.",
    "[3] 周超, 王呈, 夏源, 等. 面向工业运维人机对话的意图和语义槽联合识别算法[J]. 计算机应用研究, 2024, 41(12): 3645-3650.",
    "[4] 任元凯, 谢振平. 大语言模型领域意图的精准性增强方法[J]. 计算机应用研究, 2024, 41(10): 2893-2899.",
    "[5] 李子骏, 肖辉, 李雪峰. 面向知识密集型任务的检索增强生成技术综述[J]. 微电子学与计算机, 2025, 42(10): 48-65.",
    "[6] 曹荣荣, 柳林, 于艳东, 等. 融合知识图谱的大语言模型研究综述[J]. 计算机应用研究, 2025, 42(8): 2255-2266.",
    "[7] GAO J, GALLEY M, LI L. Neural approaches to conversational AI[J]. Foundations and Trends in Information Retrieval, 2019, 13(2-3): 127-298.",
    "[8] LEWIS P, PEREZ E, PIKTUS A, et al. Retrieval-augmented generation for knowledge-intensive NLP tasks[C]//Advances in Neural Information Processing Systems. 2020.",
    "[9] KARPUKHIN V, OGUZ B, MIN S, et al. Dense passage retrieval for open-domain question answering[C]//Proceedings of the 2020 Conference on Empirical Methods in Natural Language Processing. 2020: 6769-6781.",
    "[10] LOUVAN S, MAGNINI B. Recent neural methods on slot filling and intent classification for task-oriented dialogue systems: a survey[C]//Proceedings of the 28th International Conference on Computational Linguistics. 2020: 480-496.",
    "[11] GUU K, LEE K, TUNG Z, et al. REALM: retrieval-augmented language model pre-training[C]//Proceedings of the 37th International Conference on Machine Learning. 2020: 3929-3938.",
    "[12] CHEN M, LIU R, SHEN L, et al. The JDDC corpus: a large-scale multi-turn Chinese dialogue dataset for e-commerce customer service[C]//Proceedings of the 12th Language Resources and Evaluation Conference. 2020: 459-466.",
    "[13] DENG Y, ZHANG W, YU Q, et al. Product question answering in e-commerce: a survey[C]//Proceedings of the 61st Annual Meeting of the Association for Computational Linguistics. 2023: 11951-11964.",
    "[14] ZHU Q, ZHANG Z, FANG Y, et al. CrossWOZ: a large-scale Chinese cross-domain task-oriented dialogue dataset[J]. Transactions of the Association for Computational Linguistics, 2020, 8: 281-295.",
]

PROJECT_SCHEDULE_ROWS = [
    [
        "第8周",
        "开题材料完善",
        "完成开题报告扩写、图表重绘和参考文献整理，统一系统技术方案与答辩表述。",
        "形成可提交的开题报告终稿和答辩基础材料",
    ],
    [
        "第9周",
        "开题答辩",
        "完成开题陈述、技术路线说明和现场问答，根据教师意见明确后续实现重点。",
        "顺利完成开题答辩并进入系统实现阶段",
    ],
    [
        "第10周",
        "系统设计与核心实现",
        "推进数据库表设计、后端接口、前端页面及检索增强链路的原型开发与联调。",
        "形成可运行的系统原型和核心功能链路",
    ],
    [
        "第11周",
        "系统完善与结项准备",
        "完善订单上下文融合、多轮会话、知识库管理、测试用例和展示脚本，整理结项文档。",
        "完成结项前的系统优化、自测和材料收口",
    ],
    [
        "第12周",
        "结项答辩",
        "提交结项材料并完成系统演示、功能说明、结果汇报和现场答辩。",
        "顺利完成结项答辩与课程验收",
    ],
]


def ensure_dirs() -> None:
    DOC_OUT.mkdir(parents=True, exist_ok=True)
    TMP_DIR.mkdir(parents=True, exist_ok=True)
    FIG_DIR.mkdir(parents=True, exist_ok=True)


def set_run_font(
    run,
    name: str,
    size_pt: float,
    *,
    bold: bool = False,
    italic: bool = False,
    underline: bool = False,
    color: RGBColor | None = None,
) -> None:
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:eastAsia"), name)
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    if color is not None:
        run.font.color.rgb = color


def clear_runs(paragraph) -> None:
    p = paragraph._element
    for child in list(p):
        if child.tag != qn("w:pPr"):
            p.remove(child)


def delete_paragraph(paragraph) -> None:
    p = paragraph._element
    parent = p.getparent()
    if parent is not None:
        parent.remove(p)


def set_table_borders(table) -> None:
    tbl_pr = table._tbl.tblPr
    borders = tbl_pr.first_child_found_in("w:tblBorders")
    if borders is None:
        borders = OxmlElement("w:tblBorders")
        tbl_pr.append(borders)

    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        edge_el = borders.find(qn(f"w:{edge}"))
        if edge_el is None:
            edge_el = OxmlElement(f"w:{edge}")
            borders.append(edge_el)
        edge_el.set(qn("w:val"), "single")
        edge_el.set(qn("w:sz"), "8")
        edge_el.set(qn("w:space"), "0")
        edge_el.set(qn("w:color"), "000000")


def set_paragraph_format(paragraph, *, first_line_indent: bool = True) -> None:
    fmt = paragraph.paragraph_format
    fmt.line_spacing = 1.5
    fmt.space_before = Pt(0)
    fmt.space_after = Pt(0)
    fmt.left_indent = Pt(0)
    fmt.right_indent = Pt(0)
    fmt.first_line_indent = Pt(21) if first_line_indent else Pt(0)


def set_caption_format(paragraph) -> None:
    fmt = paragraph.paragraph_format
    fmt.line_spacing = 1.2
    fmt.space_before = Pt(3)
    fmt.space_after = Pt(3)
    fmt.left_indent = Pt(0)
    fmt.right_indent = Pt(0)
    fmt.first_line_indent = Pt(0)


def set_reference_format(paragraph) -> None:
    fmt = paragraph.paragraph_format
    fmt.line_spacing = 1.25
    fmt.space_before = Pt(0)
    fmt.space_after = Pt(0)
    fmt.left_indent = Pt(21)
    fmt.right_indent = Pt(0)
    fmt.first_line_indent = Pt(-21)


def load_font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont:
    for path in FONT_CANDIDATES:
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def text_size(draw: ImageDraw.ImageDraw, text: str, font) -> tuple[int, int]:
    bbox = draw.multiline_textbbox((0, 0), text, font=font, align="center", spacing=8)
    return bbox[2] - bbox[0], bbox[3] - bbox[1]


def wrap_text(text: str, width: int) -> str:
    lines = []
    for part in text.split("\n"):
        lines.extend(textwrap.wrap(part, width=width, break_long_words=False, replace_whitespace=False) or [""])
    return "\n".join(lines)


def draw_box(
    draw,
    box,
    text: str,
    *,
    fill: str,
    outline: str,
    font,
    text_fill: str = "#1f1f1f",
    auto_wrap: bool = False,
    wrap_width: int = 13,
) -> None:
    x1, y1, x2, y2 = box
    draw.rounded_rectangle(box, radius=20, fill=fill, outline=outline, width=4)
    wrapped = wrap_text(text, wrap_width) if auto_wrap else text
    tw, th = text_size(draw, wrapped, font)
    draw.multiline_text(
        ((x1 + x2 - tw) / 2, (y1 + y2 - th) / 2),
        wrapped,
        font=font,
        fill=text_fill,
        align="center",
        spacing=10,
    )


def draw_arrow(draw, start: tuple[int, int], end: tuple[int, int], *, color: str = "#36578c", width: int = 6) -> None:
    draw.line([start, end], fill=color, width=width)
    dx = end[0] - start[0]
    dy = end[1] - start[1]
    if dx == 0 and dy == 0:
        return
    if abs(dx) >= abs(dy):
        direction = 1 if dx > 0 else -1
        p1 = (end[0] - 18 * direction, end[1] - 10)
        p2 = (end[0] - 18 * direction, end[1] + 10)
    else:
        direction = 1 if dy > 0 else -1
        p1 = (end[0] - 10, end[1] - 18 * direction)
        p2 = (end[0] + 10, end[1] - 18 * direction)
    draw.polygon([end, p1, p2], fill=color)


def create_architecture_figure(path: Path) -> None:
    image = Image.new("RGB", (2600, 1500), "white")
    draw = ImageDraw.Draw(image)
    title_font = load_font(46, bold=True)
    box_font = load_font(40)
    note_font = load_font(30)

    title = "电商退换货智能客服系统总体架构"
    tw, th = text_size(draw, title, title_font)
    draw.text(((2600 - tw) / 2, 55), title, font=title_font, fill="#1d1d1d")

    user_box = (120, 620, 420, 840)
    vue_box = (540, 560, 980, 900)
    spring_box = (1120, 500, 1620, 960)
    mysql_box = (1820, 300, 2440, 510)
    lang_box = (1820, 620, 2440, 830)
    llm_box = (1820, 940, 2440, 1150)

    draw_box(draw, user_box, "用户端\n咨询用户 / 管理员", fill="#f4f4f4", outline="#666666", font=box_font)
    draw_box(draw, vue_box, "Vue 3 前端\n咨询工作台\n知识库管理页面", fill="#e6f2ff", outline="#3e7cb1", font=box_font)
    draw_box(
        draw,
        spring_box,
        "Spring Boot 业务核心\nREST 接口调度\n业务规则校验\n会话管理\n订单上下文服务",
        fill="#dce7f7",
        outline="#36578c",
        font=box_font,
    )
    draw_box(draw, mysql_box, "MySQL 数据层\n订单数据\nFAQ/规则文档\n会话记录与日志", fill="#f7f7f7", outline="#666666", font=box_font)
    draw_box(
        draw,
        lang_box,
        "LangChain4j AI 增强层\n提示词编排\n上下文记忆\n检索增强调用",
        fill="#e8f6ef",
        outline="#2f7d4a",
        font=box_font,
    )
    draw_box(draw, llm_box, "大模型服务\n语义理解\n自然语言生成", fill="#fff1da", outline="#b07a1d", font=box_font)

    draw_arrow(draw, (420, 730), (540, 730))
    draw_arrow(draw, (980, 730), (1120, 730))
    draw_arrow(draw, (1620, 610), (1820, 405))
    draw_arrow(draw, (1620, 725), (1820, 725))
    draw_arrow(draw, (2130, 830), (2130, 940))
    draw_arrow(draw, (1820, 450), (1620, 650))
    draw_arrow(draw, (1820, 725), (1620, 840))

    note = "设计原则：Spring Boot 负责业务核心，LangChain4j 负责 AI 增强，不以 AI 框架替代业务系统。"
    nw, nh = text_size(draw, note, note_font)
    draw.rounded_rectangle((180, 1300, 2420, 1385), radius=18, fill="#fafafa", outline="#9f9f9f", width=3)
    draw.text(((2600 - nw) / 2, 1328), note, font=note_font, fill="#404040")

    image.save(path, dpi=(300, 300))


def create_flow_figure(path: Path) -> None:
    image = Image.new("RGB", (2600, 1700), "white")
    draw = ImageDraw.Draw(image)
    title_font = load_font(46, bold=True)
    box_font = load_font(34)
    num_font = load_font(30, bold=True)
    note_font = load_font(26)

    title = "系统技术路线流程图"
    tw, th = text_size(draw, title, title_font)
    draw.text(((2600 - tw) / 2, 55), title, font=title_font, fill="#1d1d1d")

    steps = [
        "用户提交售后问题",
        "Spring Boot 接口接收请求",
        "读取订单上下文并完成规则校验",
        "检索 MySQL 中的 FAQ 与规则文档",
        "LangChain4j 组织提示词与会话记忆",
        "调用大模型生成答复并进行规则补全",
        "返回答案、依据与处理建议",
    ]
    colors = [
        ("#f4f4f4", "#666666"),
        ("#e6f2ff", "#3e7cb1"),
        ("#dce7f7", "#36578c"),
        ("#eef7ff", "#3e7cb1"),
        ("#e8f6ef", "#2f7d4a"),
        ("#fff1da", "#b07a1d"),
        ("#f4f4f4", "#666666"),
    ]

    top = 180
    height = 130
    gap = 55
    x1, x2 = 520, 2140

    for idx, step in enumerate(steps):
        y1 = top + idx * (height + gap)
        y2 = y1 + height
        draw_box(draw, (x1, y1, x2, y2), step, fill=colors[idx][0], outline=colors[idx][1], font=box_font)
        circle_x = 355
        circle_y = y1 + height // 2
        draw.ellipse((300, circle_y - 40, 410, circle_y + 40), fill=colors[idx][0], outline=colors[idx][1], width=4)
        num_text = str(idx + 1)
        nw, nh = text_size(draw, num_text, num_font)
        draw.text((circle_x - nw / 2, circle_y - nh / 2 - 4), num_text, font=num_font, fill=colors[idx][1])
        if idx < len(steps) - 1:
            draw_arrow(draw, ((x1 + x2) // 2, y2), ((x1 + x2) // 2, y2 + gap))

    note = "技术原则：先由业务服务完成规则判断与上下文读取，再由 LangChain4j 和大模型完成增强生成，最后返回可解释结果。"
    nw, nh = text_size(draw, note, note_font)
    draw.rounded_rectangle((220, 1580, 2380, 1660), radius=18, fill="#fafafa", outline="#9f9f9f", width=3)
    draw.text(((2600 - nw) / 2, 1605), note, font=note_font, fill="#404040")

    image.save(path, dpi=(300, 300))


def remove_figure_note_boxes() -> None:
    if ARCH_FIG.exists():
        image = Image.open(ARCH_FIG).convert("RGB")
        draw = ImageDraw.Draw(image)
        draw.rectangle((80, 770, 1570, 900), fill="white")
        image.save(ARCH_FIG)

    if FLOW_FIG.exists():
        image = Image.open(FLOW_FIG).convert("RGB")
        draw = ImageDraw.Draw(image)
        draw.rectangle((40, 1288, 985, 1498), fill="white")
        image.save(FLOW_FIG)


def trim_figure_whitespace(path: Path, *, threshold: int = FIGURE_BG_THRESHOLD, padding: int = FIGURE_CROP_PADDING) -> None:
    image = Image.open(path).convert("RGB")
    mask = image.convert("L").point(lambda value: 255 if value < threshold else 0, mode="1")
    bbox = mask.getbbox()
    if bbox is None:
        return

    left, top, right, bottom = bbox
    cropped = image.crop(
        (
            max(0, left - padding),
            max(0, top - padding),
            min(image.width, right + padding),
            min(image.height, bottom + padding),
        )
    )
    dpi = image.info.get("dpi", (300, 300))
    cropped.save(path, dpi=dpi)


def prepare_figure_assets() -> None:
    if GPT_ARCH_FIG.exists():
        shutil.copy2(GPT_ARCH_FIG, ARCH_FIG)
    else:
        create_architecture_figure(ARCH_FIG)

    if GPT_FLOW_FIG.exists():
        shutil.copy2(GPT_FLOW_FIG, FLOW_FIG)
    else:
        create_flow_figure(FLOW_FIG)

    remove_figure_note_boxes()
    trim_figure_whitespace(ARCH_FIG)
    trim_figure_whitespace(FLOW_FIG)


def convert_template_to_docx() -> Path:
    shutil.copy2(TEMPLATE_DOC, TMP_TEMPLATE_DOC)
    subprocess.run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "docx",
            "--outdir",
            str(TMP_DIR),
            str(TMP_TEMPLATE_DOC),
        ],
        cwd=str(ROOT),
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )
    if not TMP_TEMPLATE_DOCX.exists():
        raise FileNotFoundError(f"模板转换失败: {TMP_TEMPLATE_DOCX}")
    return TMP_TEMPLATE_DOCX


def export_pdf(input_path: Path, out_dir: Path) -> Path:
    subprocess.run(
        [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(input_path),
        ],
        cwd=str(ROOT),
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )
    pdf_path = out_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF 导出失败: {pdf_path}")
    return pdf_path


def iter_taskbook_lines() -> list[str]:
    if not TASKBOOK_PPTX.exists():
        raise FileNotFoundError(f"未找到任务书文件: {TASKBOOK_PPTX}")
    prs = Presentation(str(TASKBOOK_PPTX))
    for slide in prs.slides:
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for line in shape.text.replace("\xa0", " ").splitlines():
                    line = line.strip()
                    if line:
                        lines.append(line)
        if "课程时间内容安排" in lines:
            return [line for line in lines if line != "课程时间内容安排"]
    raise ValueError("未在任务书中找到课程时间内容安排页")


def normalize_date_text(date_text: str) -> str:
    match = re.search(r"(\d+)月(\d+)日", date_text)
    if not match:
        return date_text
    month, day = match.groups()
    return f"2026年{month}月{day}日"


def build_schedule_table() -> list[list[str]]:
    rows = [["时间节点", "任务书安排", "本项目工作内容", "阶段目标"]]
    rows.extend(PROJECT_SCHEDULE_ROWS)
    return rows


def extract_heading_pages(pdf_path: Path) -> dict[str, int]:
    doc = fitz.open(str(pdf_path.resolve()))
    page_map: dict[str, int] = {}
    try:
        for _, _, needle in TOC_ENTRY_SPECS:
            for page_index in range(doc.page_count):
                text = doc.load_page(page_index).get_text("text")
                if needle in text:
                    page_map[needle] = page_index + 1
                    break

        missing = [needle for _, _, needle in TOC_ENTRY_SPECS if needle not in page_map]
        if missing:
            raise ValueError(f"目录页码提取失败: {missing}")
        return page_map
    finally:
        doc.close()


def build_toc_entries() -> list[tuple[int, str, int]]:
    pass1_pdf = export_pdf(TMP_PASS1_DOCX, TMP_DIR)
    page_map = extract_heading_pages(pass1_pdf)
    body_start_page = page_map[TOC_ENTRY_SPECS[0][2]]
    return [
        (level, text, page_map[needle] - body_start_page + 1)
        for level, text, needle in TOC_ENTRY_SPECS
    ]


def fill_cover(doc: Document) -> None:
    cover_map = {
        "题　　目：": (TITLE, 10.5),
        "学　　院：": (ACADEMY, 14),
        "专　　业：": (MAJOR, 14),
        "班　　级：": (CLASS_NAME, 14),
        "学　　号：": (STUDENT_ID, 14),
        "姓　　名：": (STUDENT_NAME, 14),
        "指导教师：": (TEACHER, 14),
        "开题日期：": (OPENING_DATE, 14),
    }

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        for label, (value, size_pt) in cover_map.items():
            if text.startswith(label):
                clear_runs(paragraph)
                label_run = paragraph.add_run(label)
                set_run_font(label_run, "楷体_GB2312", 14)
                underlined_run = paragraph.add_run(value + "　" * 12)
                set_run_font(underlined_run, "楷体_GB2312", size_pt, underline=True)
                break


def strip_template_body(doc: Document):
    start_idx = None
    audit_idx = None
    paragraphs = list(doc.paragraphs)
    for idx, paragraph in enumerate(paragraphs):
        text = paragraph.text.strip()
        if start_idx is None and text.startswith(SECTION_TITLES["background"]):
            start_idx = idx
        if audit_idx is None and text.startswith("指导教师审核意见"):
            audit_idx = idx
            break

    if start_idx is None or audit_idx is None:
        raise ValueError("未找到模板正文或审核意见位置")

    for paragraph in paragraphs[start_idx:audit_idx]:
        delete_paragraph(paragraph)

    for paragraph in doc.paragraphs:
        if paragraph.text.strip().startswith("指导教师审核意见"):
            return paragraph
    raise ValueError("删除正文后未找到审核意见位置")


def insert_paragraph_before(anchor, text: str = ""):
    return anchor.insert_paragraph_before(text)


def set_section_type(sect_pr, value: str) -> None:
    type_el = sect_pr.find(qn("w:type"))
    if type_el is None:
        type_el = OxmlElement("w:type")
        sect_pr.insert(0, type_el)
    type_el.set(qn("w:val"), value)


def insert_section_break_before(anchor, *, section_type: str = "nextPage") -> None:
    paragraph = insert_paragraph_before(anchor)
    paragraph.paragraph_format.space_before = Pt(0)
    paragraph.paragraph_format.space_after = Pt(0)
    sect_pr = deepcopy(anchor.part.document._element.body.sectPr)
    set_section_type(sect_pr, section_type)
    paragraph._p.get_or_add_pPr().append(sect_pr)


def set_page_number_start(section, start: int) -> None:
    pg_num_type = section._sectPr.find(qn("w:pgNumType"))
    if pg_num_type is None:
        pg_num_type = OxmlElement("w:pgNumType")
        section._sectPr.append(pg_num_type)
    pg_num_type.set(qn("w:start"), str(start))


def append_page_field(run) -> None:
    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")

    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = " PAGE "

    fld_sep = OxmlElement("w:fldChar")
    fld_sep.set(qn("w:fldCharType"), "separate")

    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")

    run._r.extend([fld_begin, instr, fld_sep, fld_end])


def configure_body_page_numbers(doc: Document, *, start: int = BODY_PAGE_START) -> None:
    if len(doc.sections) < 2:
        raise ValueError("未成功创建正文节")

    section = doc.sections[1]
    section.footer.is_linked_to_previous = False
    section.footer_distance = Mm(8)
    set_page_number_start(section, start)

    footer = section.footer
    paragraph = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    clear_runs(paragraph)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.paragraph_format.space_before = Pt(0)
    paragraph.paragraph_format.space_after = Pt(0)
    paragraph.paragraph_format.line_spacing = 1.0
    paragraph.paragraph_format.first_line_indent = Pt(0)

    run = paragraph.add_run()
    set_run_font(run, "宋体", 10.5, color=BLACK)
    append_page_field(run)


def enable_update_fields_on_open(doc: Document) -> None:
    settings = doc.settings.element
    update_fields = settings.find(qn("w:updateFields"))
    if update_fields is None:
        update_fields = OxmlElement("w:updateFields")
        settings.append(update_fields)
    update_fields.set(qn("w:val"), "true")


def add_section_heading(anchor, text: str) -> None:
    paragraph = insert_paragraph_before(anchor)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
    set_paragraph_format(paragraph, first_line_indent=False)
    run = paragraph.add_run(text)
    set_run_font(run, "楷体_GB2312", 14, bold=True, color=BLACK)


def add_subheading(anchor, text: str) -> None:
    paragraph = insert_paragraph_before(anchor)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
    set_paragraph_format(paragraph, first_line_indent=False)
    run = paragraph.add_run(text)
    set_run_font(run, "宋体", 12, bold=True, color=BLACK)


def add_body_paragraph(anchor, text: str) -> None:
    paragraph = insert_paragraph_before(anchor)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    set_paragraph_format(paragraph, first_line_indent=True)
    run = paragraph.add_run(text)
    set_run_font(run, "宋体", 10.5, color=BLACK)


def add_caption(anchor, text: str) -> None:
    paragraph = insert_paragraph_before(anchor)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    set_caption_format(paragraph)
    run = paragraph.add_run(text)
    set_run_font(run, "宋体", 10.5, color=BLACK)


def add_reference(anchor, text: str) -> None:
    paragraph = insert_paragraph_before(anchor)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
    set_reference_format(paragraph)
    run = paragraph.add_run(text)
    set_run_font(run, "宋体", 10.5, color=BLACK)


def add_figure(anchor, image_path: Path, caption: str, *, width_mm: float = 145) -> None:
    paragraph = insert_paragraph_before(anchor)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraph.paragraph_format.space_before = Pt(2)
    paragraph.paragraph_format.space_after = Pt(2)
    paragraph.add_run().add_picture(str(image_path), width=Mm(width_mm))
    add_caption(anchor, caption)


def add_table(
    doc: Document,
    anchor,
    rows: list[list[str]],
    widths_mm: list[float],
    *,
    body_font_size_pt: float = 10.5,
    header_font_size_pt: float = 10.5,
) -> None:
    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
    anchor._p.addprevious(table._element)
    table.autofit = False
    set_table_borders(table)

    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.rows[row_idx].cells[col_idx]
            cell.width = Mm(widths_mm[col_idx])
            cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            paragraph = cell.paragraphs[0]
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER if row_idx == 0 else WD_ALIGN_PARAGRAPH.LEFT
            set_paragraph_format(paragraph, first_line_indent=False)
            paragraph.paragraph_format.line_spacing = 1.2
            clear_runs(paragraph)
            run = paragraph.add_run(cell_text)
            set_run_font(
                run,
                "宋体",
                header_font_size_pt if row_idx == 0 else body_font_size_pt,
                bold=(row_idx == 0),
                color=BLACK,
            )

    insert_paragraph_before(anchor)


def add_page_break_before(anchor) -> None:
    paragraph = insert_paragraph_before(anchor)
    run = paragraph.add_run()
    run.add_break(WD_BREAK.PAGE)


def add_toc_page(anchor, entries: list[tuple[int, str, int]]) -> None:
    title = insert_paragraph_before(anchor)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.paragraph_format.space_before = Pt(0)
    title.paragraph_format.space_after = Pt(12)
    title.paragraph_format.line_spacing = 1.2
    title_run = title.add_run("目录")
    set_run_font(title_run, "黑体", 16, bold=True, color=BLACK)

    for level, text, page in entries:
        paragraph = insert_paragraph_before(anchor)
        paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
        fmt = paragraph.paragraph_format
        fmt.space_before = Pt(0)
        fmt.space_after = Pt(2)
        fmt.line_spacing = 1.25
        fmt.first_line_indent = Pt(0)
        fmt.left_indent = Pt(0 if level == 1 else 21)
        fmt.right_indent = Pt(0)
        fmt.tab_stops.add_tab_stop(Mm(145), alignment=WD_TAB_ALIGNMENT.RIGHT, leader=WD_TAB_LEADER.DOTS)

        run = paragraph.add_run(text)
        set_run_font(run, "宋体", 12 if level == 1 else 10.5, bold=(level == 1), color=BLACK)
        tab_run = paragraph.add_run("\t")
        set_run_font(tab_run, "宋体", 10.5, color=BLACK)
        page_run = paragraph.add_run(str(page))
        set_run_font(page_run, "宋体", 10.5, color=BLACK)


def compose_report(doc: Document, anchor) -> None:
    add_section_heading(anchor, SECTION_TITLES["background"])
    add_subheading(anchor, "1.1 选题背景")
    for text in BACKGROUND_PARAGRAPHS:
        add_body_paragraph(anchor, text)

    add_subheading(anchor, "1.2 研究意义")
    for text in SIGNIFICANCE_PARAGRAPHS:
        add_body_paragraph(anchor, text)

    add_subheading(anchor, "1.3 国内外研究现状")
    for text in STATE_OF_ART_PARAGRAPHS:
        add_body_paragraph(anchor, text)

    add_subheading(anchor, "1.4 发展趋势")
    for text in TREND_PARAGRAPHS:
        add_body_paragraph(anchor, text)

    add_section_heading(anchor, SECTION_TITLES["content"])
    add_subheading(anchor, "2.1 研究的基本内容")
    for text in CONTENT_INTRO_PARAGRAPHS:
        add_body_paragraph(anchor, text)
    add_caption(anchor, "表2-1 研究内容与系统模块对应表")
    add_table(doc, anchor, RESEARCH_TABLE, [30, 40, 50, 35])
    add_body_paragraph(anchor, "从研究路径上看，以上五项内容彼此并非孤立模块，而是共同支撑“业务核心 + AI 增强”的整体系统方案。前端页面负责问题输入与结果展示，后端服务负责业务流程控制和规则约束，数据层沉淀订单和知识文档，LangChain4j 则作为增强层提高语义理解和答案组织能力，从而保证系统既具备完整工程结构，又具备明显的领域智能特征。")

    add_subheading(anchor, "2.2 拟解决的主要问题")
    for text in PROBLEM_PARAGRAPHS:
        add_body_paragraph(anchor, text)

    add_section_heading(anchor, SECTION_TITLES["method"])
    add_subheading(anchor, "3.1 研究方法")
    for text in METHOD_PARAGRAPHS:
        add_body_paragraph(anchor, text)

    add_subheading(anchor, "3.2 技术路线")
    for text in TECH_ROUTE_PARAGRAPHS:
        add_body_paragraph(anchor, text)
    add_figure(anchor, ARCH_FIG, "图3-1 系统总体架构图")
    add_body_paragraph(anchor, "图 3-1 从系统组成角度展示了本课题的总体结构：Vue 3 负责页面交互，Spring Boot 负责业务控制与接口服务，MySQL 负责核心数据持久化，LangChain4j 嵌入后端服务中实现检索增强和大模型编排，最终由大模型服务返回自然语言答复。该结构强调业务规则和系统控制权必须掌握在后端服务层，而 AI 相关能力以增强组件的方式嵌入其中。")
    add_body_paragraph(anchor, "从模块协同关系看，上述结构并不是把前端、后端、数据库和 AI 组件简单并列堆叠，而是围绕退换货业务形成统一闭环：前端负责承接用户输入，后端负责规则判断与接口调度，数据层负责提供订单和知识依据，AI 增强层则在此基础上完成上下文组织与答案表达。为了进一步说明这一闭环在运行阶段的先后顺序，下面结合图 3-2 对系统技术路线作进一步展开。")
    add_figure(anchor, FLOW_FIG, "图3-2 技术路线流程图", width_mm=122)
    add_body_paragraph(anchor, "图 3-2 从运行链路角度说明了系统的技术路线。用户问题首先进入 Spring Boot 接口层，由业务服务读取订单信息与会话状态，再调用知识检索和 LangChain4j 编排链路组织提示词与上下文，最终在大模型生成结果基础上输出附带依据和处理建议的答复。该流程既体现出检索增强的技术特点，也体现出课程项目对业务可控性的要求。")

    add_subheading(anchor, "3.3 研究难点")
    for text in DIFFICULTY_PARAGRAPHS:
        add_body_paragraph(anchor, text)

    add_subheading(anchor, "3.4 预期达到的目标")
    for text in GOAL_PARAGRAPHS:
        add_body_paragraph(anchor, text)

    add_section_heading(anchor, SECTION_TITLES["schedule"])
    add_body_paragraph(anchor, "结合课程当前教学进度与本课题的实施安排，本文将第8周至第12周的关键节点与项目推进目标对应起来，形成如下详细工作进度表。该表覆盖开题材料完善、开题答辩、系统实现、结项准备和结项答辩等主要阶段。")
    add_caption(anchor, "表4-1 论文详细工作进度和安排")
    add_table(doc, anchor, build_schedule_table(), [20, 28, 60, 46], body_font_size_pt=9.5)
    for text in SCHEDULE_TAIL_PARAGRAPHS:
        add_body_paragraph(anchor, text)
    for text in TEAM_WORK_INTRO_PARAGRAPHS:
        add_body_paragraph(anchor, text)
    add_caption(anchor, "表4-2 组长具体分工说明表")
    add_table(doc, anchor, TEAM_WORK_TABLE, [25, 35, 48, 46], body_font_size_pt=8, header_font_size_pt=9)

    add_section_heading(anchor, SECTION_TITLES["refs"])
    for ref in REFERENCES:
        add_reference(anchor, ref)


def build_report() -> None:
    prepare_figure_assets()

    template_docx = convert_template_to_docx()
    doc = Document(template_docx)
    section = doc.sections[0]
    section.page_height = Mm(297)
    section.page_width = Mm(210)
    section.top_margin = Mm(25.4)
    section.bottom_margin = Mm(25.4)
    section.left_margin = Mm(31.75)
    section.right_margin = Mm(31.75)

    fill_cover(doc)
    audit_para = strip_template_body(doc)
    compose_report(doc, audit_para)
    add_page_break_before(audit_para)
    doc.save(TMP_PASS1_DOCX)

    toc_entries = build_toc_entries()

    final_doc = Document(template_docx)
    final_section = final_doc.sections[0]
    final_section.page_height = Mm(297)
    final_section.page_width = Mm(210)
    final_section.top_margin = Mm(25.4)
    final_section.bottom_margin = Mm(25.4)
    final_section.left_margin = Mm(31.75)
    final_section.right_margin = Mm(31.75)

    fill_cover(final_doc)
    final_audit_para = strip_template_body(final_doc)
    add_toc_page(final_audit_para, toc_entries)
    add_page_break_before(final_audit_para)
    insert_section_break_before(final_audit_para)
    configure_body_page_numbers(final_doc)
    enable_update_fields_on_open(final_doc)
    compose_report(final_doc, final_audit_para)
    add_page_break_before(final_audit_para)

    final_doc.save(TMP_REPORT_DOCX)
    shutil.copy2(TMP_REPORT_DOCX, FINAL_DOCX)

    pdf_path = export_pdf(TMP_REPORT_DOCX, TMP_DIR)
    shutil.copy2(pdf_path, FINAL_PDF)


def main() -> None:
    ensure_dirs()
    build_report()


if __name__ == "__main__":
    main()
